from pptx import Presentation
from pptx.util import Inches


ppt = Presentation()
slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Slide em brancoasda

# Adicionar título
title = slide.shapes.title
title.text = "Gráfico de Valores Gastos"

# Salvar o PowerPoint
ppt.save("dashboard.pptx")
ppt = Presentation()
slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Slide em branco
# Adicionar imagem ao slide
img_path = "newplot.png"
left = Inches(1)
top = Inches(1.5)
slide.shapes.add_picture(img_path, left, top, width=Inches(6))


# Salvar o PowerPoint
ppt.save("dashboard.pptx")

print("Apresentação criada com sucesso!")
